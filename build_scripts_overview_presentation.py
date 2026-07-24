#!/usr/bin/env python3
"""
Builds a standalone presentation, "Data Analysis Scripts Overview", covering
how the raw data is organized and how organize_hcp_data.py,
cross_sectional_analysis_v1.py / _v2.py / _batch_v2.py, and
run_fc_pipeline_v1.py / _v2.py / _batch_v2.py each work -- narrative
companion to README.md in this same folder.

Independent of rsfMRI_processing_summary(_with_anat_appendix).pptx -- this
deck is about the scripts themselves, not a specific session's processing
results. Reuses the same design tokens/helpers established across this
project's other pptx-building scripts, plus three new slide types the brief
asked for explicitly: a full-bleed TITLE slide, full-bleed SECTION divider
slides, and image-embedding for the four existing flowchart PNGs.

Standard command:
  python3 build_scripts_overview_presentation.py

Output:
  Data_Analysis_Scripts_Overview.pptx (project root, one directory up from
  this scripts folder)
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = Path(__file__).resolve().parent
DST = HERE.parent / "presentations" / "Data_Analysis_Scripts_Overview.pptx"

# ---- design tokens, matching the rest of this project's decks ----
NAVY = RGBColor(0x0B, 0x25, 0x45)
NAVY_PANEL = RGBColor(0x14, 0x3A, 0x63)
TEAL = RGBColor(0x2E, 0xC4, 0xB6)
TEAL_DARK = RGBColor(0x1B, 0x8A, 0x80)
AMBER = RGBColor(0xE8, 0x9B, 0x3C)
LIGHT_BG = RGBColor(0xF4, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE_LIGHT = RGBColor(0xD9, 0xE1, 0xE8)
SLATE_MUTED = RGBColor(0x55, 0x66, 0x77)
BODY_DARK = RGBColor(0x1B, 0x26, 0x3B)
CAPTION_MUTED = RGBColor(0x7A, 0x8B, 0x9A)
ROW_ALT = RGBColor(0xF3, 0xF7, 0xF9)

FOOTER_TEXT = "HCP-Aging AABC Release 2  ·  Data Analysis Scripts Overview"

SLIDE_W = Emu(12191695)
SLIDE_H = Emu(6858000)
MARGIN = Emu(502920)
CONTENT_W = Emu(11201400)
CONTENT_TOP = Emu(1620000)
CONTENT_BOTTOM = Emu(6360000)

# Flowchart slides use nearly the full slide -- thinner side margins and the
# picture starts right below the accent bar instead of below a subtitle line.
FLOWCHART_MARGIN = Emu(120000)
FLOWCHART_TOP = Emu(1130000)
FLOWCHART_W = SLIDE_W - 2 * FLOWCHART_MARGIN
FLOWCHART_H = Emu(6380000) - FLOWCHART_TOP

PAGE = [0]  # mutable page counter


def next_page():
    PAGE[0] += 1
    return PAGE[0]


def set_solid_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size, color, bold=False,
                 font="Calibri", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_header(slide, eyebrow, title, page_num):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Emu(1051560))
    band.fill.solid(); band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Emu(1051560), SLIDE_W, Emu(38100))
    accent.fill.solid(); accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    add_textbox(slide, MARGIN, Emu(146304), Emu(9144000), Emu(274320),
                eyebrow, 12, TEAL, bold=True)
    add_textbox(slide, MARGIN, Emu(402336), Emu(10515600), Emu(594360),
                title, 26, WHITE, bold=True)
    add_textbox(slide, Emu(11277295), Emu(146304), Emu(548640), Emu(274320),
                str(page_num), 12, SLATE_LIGHT)


def add_footer(slide, page_num):
    add_textbox(slide, MARGIN, Emu(6492240), Emu(7315200), Emu(274320),
                FOOTER_TEXT, 9, SLATE_MUTED)
    add_textbox(slide, Emu(11094415), Emu(6492240), Emu(640080), Emu(274320),
                str(page_num), 9, SLATE_MUTED)


def new_content_slide(prs, blank_layout, eyebrow, title):
    page_num = next_page()
    slide = prs.slides.add_slide(blank_layout)
    set_solid_bg(slide, LIGHT_BG)
    add_header(slide, eyebrow, title, page_num)
    add_footer(slide, page_num)
    return slide


def new_title_slide(prs, blank_layout, eyebrow, title, subtitle, meta):
    """Full-bleed navy title slide -- distinct from content slides, no header band."""
    page_num = next_page()
    slide = prs.slides.add_slide(blank_layout)
    set_solid_bg(slide, NAVY)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Emu(6700000), SLIDE_W, Emu(38100))
    accent.fill.solid(); accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    add_textbox(slide, MARGIN, Emu(2400000), Emu(11186760), Emu(320000),
                eyebrow, 14, TEAL, bold=True)
    add_textbox(slide, MARGIN, Emu(2760000), Emu(11186760), Emu(1200000),
                title, 40, WHITE, bold=True)
    add_textbox(slide, MARGIN, Emu(3900000), Emu(11186760), Emu(500000),
                subtitle, 17, SLATE_LIGHT)
    add_textbox(slide, MARGIN, Emu(6200000), Emu(11186760), Emu(300000),
                meta, 11, SLATE_MUTED)
    return slide


def new_section_slide(prs, blank_layout, section_num, section_title, description):
    """Full-bleed teal-accented section divider -- distinct from content slides."""
    page_num = next_page()
    slide = prs.slides.add_slide(blank_layout)
    set_solid_bg(slide, NAVY)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Emu(3100000), Emu(200000), Emu(900000))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    add_textbox(slide, Emu(900000), Emu(2950000), Emu(4000000), Emu(320000),
                f"SECTION {section_num:02d}", 15, TEAL, bold=True)
    add_textbox(slide, Emu(900000), Emu(3280000), Emu(10300000), Emu(900000),
                section_title, 34, WHITE, bold=True)
    add_textbox(slide, Emu(900000), Emu(4200000), Emu(9000000), Emu(500000),
                description, 14, SLATE_LIGHT)
    add_textbox(slide, MARGIN, Emu(6492240), Emu(640080), Emu(274320),
                str(page_num), 9, RGBColor(0x8A, 0x9B, 0xB0))
    return slide


def add_table(slide, left, top, width, col_widths, rows, header_size=11.5, body_size=10,
              row_h=Emu(320000), mono_cols=None, highlight_rows=None):
    mono_cols = mono_cols or set()
    highlight_rows = highlight_rows or set()
    n_rows = len(rows)
    n_cols = len(rows[0])
    height = row_h * n_rows
    gshape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gshape.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
    for ri, row in enumerate(rows):
        table.rows[ri].height = row_h
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.margin_left = Emu(91440)
            cell.margin_right = Emu(45720)
            cell.margin_top = Emu(22860)
            cell.margin_bottom = Emu(22860)
            cell.vertical_anchor = 3  # MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = val
            is_header = ri == 0
            is_mono = ci in mono_cols and not is_header
            run.font.name = "Menlo" if is_mono else "Calibri"
            run.font.size = Pt(header_size if is_header else body_size)
            run.font.bold = is_header
            if is_header:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
                run.font.color.rgb = WHITE
            elif (ri - 1) in highlight_rows:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xE4, 0xF7, 0xF5)
                run.font.color.rgb = TEAL_DARK if is_mono else BODY_DARK
                run.font.bold = is_mono
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ROW_ALT if (ri % 2 == 0) else WHITE
                run.font.color.rgb = BODY_DARK
    return gshape


def add_callout_panel(slide, left, top, w, h, header, bullets, header_color=TEAL, body_size=11.5):
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    panel.fill.solid(); panel.fill.fore_color.rgb = NAVY_PANEL
    panel.line.fill.background()
    pad = Emu(228600)
    add_textbox(slide, left + pad, top + Emu(210312), w - Emu(457200), Emu(274320),
                header, 11, header_color, bold=True)
    body = "\n".join(f"▸  {b}" for b in bullets)
    add_textbox(slide, left + pad, top + Emu(575872), w - Emu(411480), h - Emu(700000),
                body, body_size, WHITE)


def add_bullets_panel(slide, left, top, w, h, header, bullets, accent=TEAL_DARK, body_size=12):
    add_textbox(slide, left, top, w, Emu(320000), header, 15, NAVY, bold=True)
    body = "\n".join(f"▸  {b}" for b in bullets)
    add_textbox(slide, left, top + Emu(400000), w, h - Emu(400000), body, body_size, BODY_DARK)


def add_picture_fit(slide, img_path, area_left, area_top, area_w, area_h):
    """Fit an image into (area_w, area_h), preserving aspect ratio, centered."""
    with Image.open(img_path) as im:
        iw, ih = im.size
    img_ratio = iw / ih
    area_ratio = area_w / area_h
    if img_ratio > area_ratio:
        w = area_w
        h = int(w / img_ratio)
    else:
        h = area_h
        w = int(h * img_ratio)
    left = area_left + (area_w - w) // 2
    top = area_top + (area_h - h) // 2
    slide.shapes.add_picture(str(img_path), left, top, width=w, height=h)


def add_tree_box(slide, left, top, w, h, text, size=11.5):
    box = slide.shapes.add_textbox(left, top, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = SLATE_LIGHT; box.line.width = Emu(9525)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(160000); tf.margin_right = Emu(120000)
    tf.margin_top = Emu(120000); tf.margin_bottom = Emu(120000)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run(); run.text = line
        run.font.name = "Menlo"; run.font.size = Pt(size); run.font.color.rgb = BODY_DARK


def add_stat_card(slide, left, top, w, h, label, value, sub):
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = SLATE_LIGHT; card.line.width = Emu(12700)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Emu(50800), h)
    accent.fill.solid(); accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    pad = Emu(180000)
    add_textbox(slide, left + pad, top + Emu(110000), w - Emu(330000), Emu(260000),
                label, 10.5, SLATE_MUTED, bold=True)
    add_textbox(slide, left + pad, top + Emu(340000), w - Emu(330000), Emu(420000),
                value, 20, NAVY, bold=True)
    add_textbox(slide, left + pad, top + Emu(760000), w - Emu(330000), Emu(400000),
                sub, 9.5, SLATE_MUTED)


# =========================================================================

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    # ================= TITLE SLIDE =================
    new_title_slide(
        prs, blank_layout,
        "HCP-AGING · AABC RELEASE 2",
        "Data Analysis Scripts — Overview",
        "How the raw data is organized, and how each script works: organize_hcp_data.py,\n"
        "cross_sectional_analysis*.py (v1/v2 + batch), and run_fc_pipeline*.py (v1/v2 + batch)",
        "data_analysis_scripts/  ·  companion to README.md and README.txt in the same folder",
    )

    # ================= CONTENTS =================
    s = new_content_slide(prs, blank_layout, "CONTENTS", "Agenda")
    agenda = [
        ("01", "Data Organization & Hierarchy", "Raw data layout, Analysed_data mirror, shared atlases"),
        ("02", "organize_hcp_data.py", "Zip → sub-/ses-/anat|func|concat, step by step"),
        ("03", "Cross-Sectional Anatomical Analysis", "v1 (thickness/myelin/volumes) → v2 (+ HCPex midbrain/basal forebrain) → batch driver"),
        ("04", "Functional Connectivity Pipeline", "v1 → v2 atlas architecture, plus a batch driver for v2"),
    ]
    y = Emu(1750000)
    for num, title, desc in agenda:
        add_textbox(s, MARGIN, y, Emu(900000), Emu(320000), num, 20, TEAL, bold=True)
        add_textbox(s, MARGIN + Emu(950000), y, Emu(9700000), Emu(320000), title, 17, NAVY, bold=True)
        add_textbox(s, MARGIN + Emu(950000), y + Emu(360000), Emu(9700000), Emu(280000), desc, 11.5, SLATE_MUTED)
        y += Emu(1100000)

    # ================= SECTION 01 =================
    new_section_slide(prs, blank_layout, 1, "Data Organization & Hierarchy",
                       "Raw data layout, the Analysed_data mirror, and the shared atlases folder")

    # ---- 01a: Raw data hierarchy ----
    s = new_content_slide(prs, blank_layout, "SECTION 01 · DATA ORGANIZATION", "Raw Data Hierarchy")
    add_textbox(s, MARGIN, Emu(1230000), CONTENT_W, Emu(274320),
                "Every script in this folder expects this layout as its raw-data root, built by organize_hcp_data.py (Section 02).",
                11, SLATE_MUTED)
    tree = (
        "Raw_Data/\n"
        "├─ sub-<subject>/                    e.g. sub-HCA6002236\n"
        "│   └─ ses-<visit>/                  e.g. ses-V3\n"
        "│       ├─ anat/                     structural (T1w/T2w) preprocessing\n"
        "│       │   ├─ T1w/                  native (acpc) anatomicals + FreeSurfer seg.\n"
        "│       │   ├─ MNINonLinear/         same, registered to MNI152\n"
        "│       │   ├─ unprocessed/          raw acquisition metadata only\n"
        "│       │   └─ ProcessingInfo/       QuNex + HCP pipeline logs/scripts\n"
        "│       ├─ func/                     resting-state fMRI (per-run + concatenated)\n"
        "│       └─ concat/                   concatenated multi-run fMRI outputs\n"
        "└─ atlases/                          shared atlas files (next slide)"
    )
    add_tree_box(s, MARGIN, Emu(1620000), CONTENT_W, Emu(4600000), tree, size=12.5)

    # ---- 01b: Analysed_data + atlases ----
    s = new_content_slide(prs, blank_layout, "SECTION 01 · DATA ORGANIZATION", "Analysed_data Mirror & Shared Atlases")
    add_textbox(s, MARGIN, Emu(1230000), CONTENT_W, Emu(274320),
                "Analysed_data/ sits next to (not inside) the raw data root, and mirrors its sub-*/ses-* layout exactly. Raw data is never written to.",
                11, SLATE_MUTED)
    tree2 = (
        "Analysed_data/\n"
        "├─ analysis_log*.json                one run-count log per pipeline, never shared\n"
        "└─ <subject>/\n"
        "    ├─ timeseries_<ses>.csv          run_fc_pipeline_v1.py's cached timeseries\n"
        "    ├─ timeseries_hcpex_<ses>.csv    run_fc_pipeline_v2.py's cached timeseries\n"
        "    └─ <session>/\n"
        "        ├─ anat/                     cross_sectional_analysis_v1.py / _v2.py\n"
        "        └─ func/                     run_fc_pipeline_v1.py / _v2.py (shared, suffix per mode)"
    )
    add_tree_box(s, MARGIN, Emu(1620000), Emu(6600000), Emu(2900000), tree2, size=11.5)

    add_textbox(s, MARGIN, Emu(4750000), Emu(6600000), Emu(300000), "Shared Atlases (Raw_Data/atlases/)", 14, NAVY, bold=True)
    rows_atlas = [
        ["File", "Used by", "Covers"],
        ["schaefer400_tianS1.dlabel.nii", "v1, cross-sec.", "400 Schaefer cortex + 16 Tian-S1 subcortex"],
        ["CIT168_prob_func2mm.nii.gz", "v1", "5 midbrain nuclei, probabilistic"],
        ["HCPex_2mm.nii + labels", "v1 (short codes), v2 + cross-sec. v2 (full names)", "426: 360 cortex + 66 sub/midbrain/BF"],
        ["Tian_Subcortex_S1_3T.nii", "v1", "Official volumetric Tian-S1"],
    ]
    add_table(s, MARGIN, Emu(5100000), Emu(6600000),
              [Emu(2700000), Emu(1800000), Emu(2100000)],
              rows_atlas, header_size=9.5, body_size=8.5, row_h=Emu(235000), mono_cols={0})

    add_callout_panel(
        s, Emu(7300000), Emu(1620000), Emu(3400000), Emu(3730000),
        "KEY INVARIANTS",
        [
            "Analysed_data/ layout mirrors Raw_Data/ 1:1 by subject/session.",
            "Every pipeline version keeps its own run-count log — never overwrites another version's.",
            "Atlas files are shared read-only inputs — no script ever writes into atlases/.",
            "A FALLBACK_ATLASES path (/Volumes/njainmpi/...) is baked into every script in case the local atlases/ folder is missing.",
        ],
        body_size=11,
    )

    # ================= SECTION 02 =================
    new_section_slide(prs, blank_layout, 2, "organize_hcp_data.py",
                       "Turning downloaded zip packages into the sub-/ses-/anat|func|concat hierarchy")

    s = new_content_slide(prs, blank_layout, "SECTION 02 · organize_hcp_data.py", "What It Does & How")
    add_bullets_panel(
        s, MARGIN, Emu(1300000), CONTENT_W, Emu(2500000),
        "Six steps, one zip at a time",
        [
            "Parse the filename — a regex pulls subject, visit, and modality "
            "(Structural/RestFmri/ConcatFmri) directly from the zip name, e.g. "
            "HCA6002236_V3_MR_StructuralRecommended....zip",
            "Route to a modality folder — Structural→anat, RestFmri→func, ConcatFmri→concat.",
            "Skip work already done — already_extracted() checks the target folder first.",
            "Extract to scratch, then merge — resolve_content_root() skips past the single "
            "wrapper folder every HCP zip extracts into; merge_tree() moves files in without "
            "ever overwriting an existing file (size mismatches go to conflicts.log).",
            "Archive the zip — moved to output_root/archive/ once merged, keeping "
            "processed and unprocessed zips separate.",
            "Report — one append-only row per zip in manifest.csv (subject, visit, "
            "modality, status, timestamp) — re-running against mixed old/new zips "
            "builds one continuous audit trail.",
        ],
        body_size=12.5,
    )
    add_textbox(s, MARGIN, Emu(4200000), CONTENT_W, Emu(280000),
                "python3 organize_hcp_data.py --input /path/to/zip_folder --output /path/to/organized_root  [--dry-run] [--subject HCA6072156]",
                11.5, TEAL_DARK, font="Menlo")

    s = new_content_slide(prs, blank_layout, "SECTION 02 · organize_hcp_data.py", "Flowchart")
    add_picture_fit(s, HERE / "organize_hcp_data_flowchart.png", FLOWCHART_MARGIN, FLOWCHART_TOP, FLOWCHART_W, FLOWCHART_H)

    # ================= SECTION 03 =================
    new_section_slide(prs, blank_layout, 3, "Cross-Sectional Anatomical Analysis",
                       "cross_sectional_analysis_v1.py, _v2.py, and a batch driver — cortical thickness, myelin, regional volumes")

    s = new_content_slide(prs, blank_layout, "SECTION 03 · CROSS-SECTIONAL ANALYSIS", "v1 — cross_sectional_analysis_v1.py")
    add_bullets_panel(
        s, MARGIN, Emu(1300000), CONTENT_W, Emu(3400000),
        "Three structural measures, per session",
        [
            "Cortical thickness & myelin (T1w/T2w ratio) — start as per-vertex CIFTI "
            "dscalar maps (corrThickness_MSMAll, MyelinMap_BC_MSMAll) with no built-in "
            "regions, so they're parcellated to Schaefer-400 by matching (structure, "
            "vertex index) between atlas and dscalar in pure Python/nibabel — no "
            "wb_command dependency, verified to match wb_command -cifti-parcellate to ~1e-8.",
            "Subcortical/regional volumes — every FreeSurfer label present in "
            "wmparc.nii.gz, volume = voxel count × voxel volume READ FROM THE FILE'S "
            "OWN AFFINE (this dataset: 0.512 mm³/voxel, 0.8mm iso — not 2mm).",
            "No atlas needed for volumes — wmparc/aparc+aseg are already FreeSurfer's "
            "own atlas-labeled segmentation output; the atlas step already happened "
            "upstream, during structural preprocessing.",
        ],
        body_size=13,
    )
    add_textbox(s, MARGIN, Emu(5000000), CONTENT_W, Emu(280000),
                "Output: Analysed_data/<subject>/<session>/anat/  ·  cortical_thickness_schaefer400.csv, "
                "myelin_schaefer400.csv, subcortical_volumes_wmparc.csv, 2 PNGs, manifest.txt",
                10.5, CAPTION_MUTED)

    s = new_content_slide(prs, blank_layout, "SECTION 03 · CROSS-SECTIONAL ANALYSIS", "v2 — Adds VTA / SN / Nucleus Basalis")
    add_textbox(s, MARGIN, Emu(1230000), CONTENT_W, Emu(500000),
                "Everything v1 does, plus volumes for 3 structures FreeSurfer's wmparc doesn't segment at all — "
                "pulled from HCPex instead (same atlas run_fc_pipeline_v2.py uses for these regions).",
                12, SLATE_MUTED)
    add_callout_panel(
        s, MARGIN, Emu(1900000), CONTENT_W, Emu(2300000),
        "THE METHODOLOGICAL PROBLEM V2 HAD TO SOLVE",
        [
            "HCPex is a template atlas fixed in standard MNI152 space — every subject's "
            "MNINonLinear/ data was already warped INTO that space.",
            "Counting voxels directly in HCPex gives every subject the IDENTICAL volume "
            "(same voxels, by construction) — zero cross-subject variance, useless for "
            "cross-sectional comparison.",
        ],
        header_color=AMBER,
        body_size=13,
    )
    rows_v2 = [
        ["Method", "How", "Varies by subject?"],
        ["Standard-space volume", "Raw voxel count in HCPex_2mm.nii (2mm grid)", "No — fixed reference only"],
        ["Native-space volume", "HCPex warped via FSL applywarp (--warp=standard2acpc_dc.nii.gz "
         "--ref=T1w_acpc_dc.nii.gz --interp=nn), counted at native voxel size", "Yes — correct for cross-sectional comparison"],
    ]
    add_table(s, MARGIN, Emu(4350000), CONTENT_W,
              [Emu(2400000), Emu(6100000), Emu(2701400)],
              rows_v2, header_size=11, body_size=10, row_h=Emu(620000), highlight_rows={1})

    s = new_content_slide(prs, blank_layout, "SECTION 03 · CROSS-SECTIONAL ANALYSIS", "Batch Driver — cross_sectional_analysis_batch_v2.py")
    add_bullets_panel(
        s, MARGIN, Emu(1300000), CONTENT_W, Emu(3000000),
        "Non-interactive driver for cross_sectional_analysis_v2.py",
        [
            "Imports cross_sectional_analysis_v2.py as a module and reuses its "
            "extraction/plotting/logging functions unchanged — only the interactive "
            "subject/session picker is replaced with a loop over every sub-*/ses-* "
            "session found under a raw data root.",
            "Work identical for every subject — FreeSurfer/HCPex LUTs, the cortex "
            "vertex LUT, and HCPex's standard-space volumes (fixed by construction) "
            "— is computed once up front instead of per session.",
            "Sessions that already have anat/manifest.txt are skipped unless --force; "
            "each session runs in its own try/except so one failure doesn't abort "
            "the batch. Shares analysis_log_anat_v2.json with the interactive script.",
        ],
        body_size=13,
    )
    add_textbox(s, MARGIN, Emu(4550000), CONTENT_W, Emu(280000),
                "python3 cross_sectional_analysis_batch_v2.py /path/to/Raw_Data  [--force] [--subjects sub-A,sub-B]",
                11.5, TEAL_DARK, font="Menlo")

    s = new_content_slide(prs, blank_layout, "SECTION 03 · CROSS-SECTIONAL ANALYSIS", "Flowchart — How Every Value Is Computed")
    add_picture_fit(s, HERE / "cross_sectional_analysis_flowchart.png", FLOWCHART_MARGIN, FLOWCHART_TOP, FLOWCHART_W, FLOWCHART_H)

    s = new_content_slide(prs, blank_layout, "SECTION 03 · CROSS-SECTIONAL ANALYSIS", "v1 vs v2 — Differences")
    rows_diff = [
        ["", "v1", "v2"],
        ["Cortical thickness / myelin", "✓ Schaefer-400", "✓ Schaefer-400 (unchanged)"],
        ["wmparc regional volumes", "✓ every FreeSurfer label", "✓ every FreeSurfer label (unchanged)"],
        ["VTA / SN / Nucleus Basalis", "— not covered", "✓ via HCPex, standard + native space"],
        ["External dependency", "nibabel only", "+ FSL applywarp on PATH"],
        ["Extra output file", "—", "midbrain_basalforebrain_volumes_hcpex.csv/.png"],
        ["Run-count log", "analysis_log_anat.json", "analysis_log_anat_v2.json (separate)"],
        ["Batch driver", "— none", "cross_sectional_analysis_batch_v2.py"],
    ]
    add_table(s, MARGIN, Emu(1620000), CONTENT_W,
              [Emu(3200000), Emu(3900000), Emu(4101400)],
              rows_diff, header_size=12, body_size=10.5, row_h=Emu(440000), highlight_rows={2})
    add_textbox(s, MARGIN, Emu(5300000), CONTENT_W, Emu(400000),
                "Shared caveat: no aseg.stats (or equivalent eTIV source) exists anywhere in this dataset's "
                "structural output — all volumes are raw mm³, not head-size-normalized.",
                10.5, CAPTION_MUTED)

    # ================= SECTION 04 =================
    new_section_slide(prs, blank_layout, 4, "Functional Connectivity Pipeline",
                       "run_fc_pipeline_v1.py → _v2.py, plus a batch driver for v2")

    s = new_content_slide(prs, blank_layout, "SECTION 04 · FC PIPELINE", "v1 — run_fc_pipeline_v1.py")
    add_bullets_panel(
        s, MARGIN, Emu(1300000), CONTENT_W, Emu(3800000),
        "The four-atlas-group interactive pipeline — 427 parcels",
        [
            "Group A (cortex) — 400 Schaefer parcels via wb_command -cifti-parcellate "
            "on the grayordinate dtseries; the only group still on the CIFTI/MSMAll "
            "path, since surface-registration precision is the reason to keep it there.",
            "Group B (subcortex) — 16 Tian-S1 parcels, mask-mean volumetric extraction "
            "from the BOLD (shares one file + method with Groups C/D, not CIFTI).",
            "Group C (midbrain) — 5 CIT168 nuclei (SNc, SNr, VTA, PBP, RN), "
            "probability-weighted volumetric extraction.",
            "Group D (basal forebrain) — 6 HCPex structures (mammillary bodies, septal "
            "nucleus, nucleus basalis ×2 hemi), mask-mean volumetric extraction.",
            "Pick any subset across the four groups (e.g. 'B1, C3, D5-6') before computing "
            "the Pearson + Fisher-z FC matrix and a labeled heatmap. Group A's atlas card "
            "shows a human-readable display label (e.g. 'A107 L SalVentAttn Medial #1') — "
            "display only; output files still use the raw Schaefer name.",
        ],
        body_size=12,
    )
    add_textbox(s, MARGIN, Emu(5600000), CONTENT_W, Emu(280000),
                "Output: Analysed_data/<subject>/<session>/func/  ·  fc_matrix_corr_<mode>.csv, "
                "fc_matrix_fisherz_<mode>.csv, region_names_<mode>.txt, fc_matrix_<mode>.png/.svg",
                10.5, CAPTION_MUTED)

    s = new_content_slide(prs, blank_layout, "SECTION 04 · FC PIPELINE", "v1 — Five Fixed Analyses (No Mode Menu)")
    add_textbox(s, MARGIN, Emu(1230000), CONTENT_W, Emu(500000),
                "After picking a subject/session and which parcels go into the \"standard\" matrix, "
                "every run always computes the same fixed five analyses — no mode menu, "
                "four of five have no choice to make anyway:",
                12, SLATE_MUTED)
    rows_v2 = [
        ["standard", "graph_vta", "graph_sn", "triangle_vta", "triangle_sn"],
        ["Any parcels, full 427-parcel atlas", "VTA + NbM(L/R) + HC(L/R) — 5 nodes, L/R kept separate",
         "SNc + SNr + NbM(L/R) + HC(L/R) — 6 nodes, L/R kept separate", "Same nodes as graph_vta, L/R averaged into one node",
         "Same idea for SN (SNc+SNr averaged too)"],
    ]
    add_table(s, MARGIN, Emu(4350000), CONTENT_W,
              [Emu(2200000), Emu(2250350), Emu(2250350), Emu(2250350), Emu(2250350)],
              rows_v2, header_size=10.5, body_size=8.2, row_h=Emu(750000))
    add_textbox(s, MARGIN, Emu(5950000), CONTENT_W, Emu(400000),
                "Output: Analysed_data/<subject>/<session>/ (no per-mode subfolder), filenames suffixed by analysis, plain CSV not .npy.",
                10.5, CAPTION_MUTED)

    s = new_content_slide(prs, blank_layout, "SECTION 04 · FC PIPELINE", "v2 — HCPex-Only, Single Atlas")
    add_textbox(s, MARGIN, Emu(1230000), CONTENT_W, Emu(500000),
                "Not just an extension of v1 — a different architecture. Every parcel, cortex included, comes from ONE atlas.",
                12, SLATE_MUTED)
    add_callout_panel(
        s, MARGIN, Emu(1850000), CONTENT_W, Emu(1900000),
        "ONE ATLAS, ONE FILE, ONE METHOD",
        [
            "HCPex: 426 regions = 360 HCP-MMP1.0 cortex + 66 subcortical/midbrain/"
            "basal-forebrain — all extracted by volumetric mask-mean from one file.",
            "Tradeoff: loses MSMAll's surface-registration precision for cortex, "
            "in exchange for one unified pipeline instead of four atlas groups.",
        ],
        body_size=12.5,
    )
    add_textbox(s, MARGIN, Emu(3950000), CONTENT_W, Emu(300000),
                "Every run always computes all 5 analyses for the chosen session:", 13, NAVY, bold=True)
    rows_v3 = [
        ["standard", "graph_vta", "graph_sn", "triangle_vta", "triangle_sn"],
        ["Any parcels, full 426-parcel atlas", "NbM↔VTA↔Hippocampus, L/R kept separate",
         "NbM↔SN↔Hippocampus, L/R kept separate", "Same 3 regions, L/R averaged into one node",
         "Same idea for SN (SNpc+SNpr averaged)"],
    ]
    add_table(s, MARGIN, Emu(4350000), CONTENT_W,
              [Emu(2200000), Emu(2250350), Emu(2250350), Emu(2250350), Emu(2250350)],
              rows_v3, header_size=10.5, body_size=8.5, row_h=Emu(750000))
    add_textbox(s, MARGIN, Emu(5950000), CONTENT_W, Emu(400000),
                "Output: same Analysed_data/<subject>/<session>/ folder as v1/v2, filenames suffixed \"_hcpex\", plain CSV not .npy.",
                10.5, CAPTION_MUTED)

    s = new_content_slide(prs, blank_layout, "SECTION 04 · FC PIPELINE", "Batch Driver — run_fc_pipeline_batch_v2.py")
    add_bullets_panel(
        s, MARGIN, Emu(1300000), CONTENT_W, Emu(3000000),
        "Non-interactive driver for run_fc_pipeline_v2.py (HCPex-only)",
        [
            "Imports run_fc_pipeline_v2.py as a module and reuses its extraction/"
            "analysis/plotting/logging functions unchanged — only the interactive "
            "subject/session/parcel pickers are replaced with a loop over every "
            "sub-*/ses-* session found under a raw data root.",
            "The \"standard\" matrix's parcel selection is made ONCE up front (same "
            "syntax as the interactive card picker — 'all', a bare group letter, or "
            "labels/ranges like 'A1, B2, B4' / 'A1-10') and applied identically to "
            "every session, since HCPex's 426 parcel names/order are fixed by the atlas.",
            "Sessions that already have a standard FC matrix are skipped unless --force; "
            "each session runs in its own try/except so one failure doesn't abort the "
            "batch. Shares analysis_log_v3.json with the interactive run_fc_pipeline_v2.py.",
        ],
        body_size=12.5,
    )
    add_textbox(s, MARGIN, Emu(4550000), CONTENT_W, Emu(280000),
                "python3 run_fc_pipeline_batch_v2.py /path/to/Raw_Data  [--force] [--parcels all] [--subjects sub-A,sub-B]",
                11.5, TEAL_DARK, font="Menlo")
    add_textbox(s, MARGIN, Emu(5000000), CONTENT_W, Emu(400000),
                "There is no batch driver for run_fc_pipeline_v1.py in this folder — the one that used to fill "
                "that role, batching the old CIT168/Tian-S1 pipeline, was removed along with run_fc_pipeline.py.",
                10.5, CAPTION_MUTED)

    s = new_content_slide(prs, blank_layout, "SECTION 04 · FC PIPELINE", "Flowchart — Correlation Estimation")
    add_picture_fit(s, HERE / "fc_correlation_estimation_flowchart.png", FLOWCHART_MARGIN, FLOWCHART_TOP, FLOWCHART_W, FLOWCHART_H)

    s = new_content_slide(prs, blank_layout, "SECTION 04 · FC PIPELINE", "Flowchart — v1 Pipeline Architecture")
    add_picture_fit(s, HERE / "fc_pipeline_v1_flowchart.png", FLOWCHART_MARGIN, FLOWCHART_TOP, FLOWCHART_W, FLOWCHART_H - Emu(320000))
    add_textbox(s, FLOWCHART_MARGIN, Emu(6110000), FLOWCHART_W, Emu(260000),
                "No editable source — recovered from this deck's own embedded media after the standalone file "
                "went missing. Its title text still reads \"Pipeline v2\", predating the 2026-07-23 rename.",
                8.5, CAPTION_MUTED, align=PP_ALIGN.CENTER)

    s = new_content_slide(prs, blank_layout, "SECTION 04 · FC PIPELINE", "v1 vs v2 — Differences")
    rows_fcdiff = [
        ["", "v1", "v2"],
        ["Cortex source", "CIFTI (Schaefer-400)", "Volumetric (HCPex)"],
        ["Subcortex source", "Volumetric (Tian-S1)", "Volumetric (HCPex)"],
        ["Midbrain source", "Volumetric (CIT168)", "Volumetric (HCPex)"],
        ["Basal forebrain source", "Volumetric (HCPex)", "Volumetric (HCPex)"],
        ["Total parcels", "427", "426"],
        ["Atlas files needed", "4", "1"],
        ["Analysis modes", "5 (fixed, no menu)", "5 (fixed, no menu)"],
        ["Output format", ".csv", ".csv"],
        ["Run-count log", "analysis_log_v2.json", "analysis_log_v3.json"],
        ["Batch driver", "— none", "run_fc_pipeline_batch_v2.py"],
    ]
    add_table(s, MARGIN, Emu(1620000), CONTENT_W,
              [Emu(3400000), Emu(3900000), Emu(3901400)],
              rows_fcdiff, header_size=11, body_size=10, row_h=Emu(430000), highlight_rows={5})

    # ================= PREREQUISITES =================
    s = new_content_slide(prs, blank_layout, "SUMMARY", "System Prerequisites — Everything Needed Before Running")
    add_textbox(s, MARGIN, Emu(1230000), CONTENT_W, Emu(274320),
                "External tools, Python libraries, atlas/lookup files, and data-format assumptions — versions shown are what's verified installed on this machine.",
                11, SLATE_MUTED)

    quad_gap = Emu(280000)
    quad_w = (CONTENT_W - quad_gap) // 2
    quad2_left = MARGIN + quad_w + quad_gap
    row1_top = Emu(1580000)
    row1_h = Emu(2150000)
    row2_top = Emu(3900000)
    row2_h = Emu(2350000)

    # ---- Quadrant 1: External command-line tools ----
    add_textbox(s, MARGIN, row1_top, quad_w, Emu(260000), "External Tools", 12.5, NAVY, bold=True)
    rows_tools = [
        ["Tool", "Version (verified)", "Required by"],
        ["Python 3", "3.14.6", "Every script"],
        ["Connectome Workbench (wb_command)", "2.2.1", "run_fc_pipeline_v1.py (Group A CIFTI)"],
        ["FSL (applywarp)", "6.0.7", "cross_sectional_analysis_v2.py (native-space warp)"],
        ["FreeSurfer", "8.1.0", "cross_sectional_analysis*.py (FreeSurferColorLUT.txt)"],
        ["Graphviz (dot)", "14.1.2", "Regenerating .dot flowcharts only — not the analyses"],
    ]
    add_table(s, MARGIN, row1_top + Emu(300000), quad_w,
              [Emu(2600000), Emu(1150000), Emu(1700700)],
              rows_tools, header_size=8.5, body_size=7.8, row_h=Emu(310000), mono_cols={0})

    # ---- Quadrant 2: Python libraries ----
    add_textbox(s, quad2_left, row1_top, quad_w, Emu(260000), "Python Libraries", 12.5, NAVY, bold=True)
    rows_py = [
        ["Package", "Version", "Used by"],
        ["numpy", "2.4.6", "cross_sectional_analysis*.py, run_fc_pipeline*.py"],
        ["nibabel", "5.4.2", "Reading CIFTI/NIfTI — same scripts"],
        ["matplotlib", "3.11.0", "All PNG/SVG plots — same scripts"],
        ["python-pptx", "1.0.2", "Every *_slides.py / build_*_presentation.py"],
        ["Pillow (PIL)", "12.2.0", "build_scripts_overview_presentation.py"],
    ]
    add_table(s, quad2_left, row1_top + Emu(300000), quad_w,
              [Emu(1750000), Emu(1050000), Emu(2650700)],
              rows_py, header_size=8.5, body_size=7.8, row_h=Emu(310000), mono_cols={0})

    # ---- Quadrant 3: Atlas & lookup files ----
    add_textbox(s, MARGIN, row2_top, quad_w, Emu(260000), "Atlas & Lookup Files (Raw_Data/atlases/)", 12.5, NAVY, bold=True)
    atlas_tree = (
        "schaefer400_tianS1.dlabel.nii          Schaefer-400 + Tian-S1 (CIFTI)\n"
        "CIT168_prob_func2mm.nii.gz + labels    5 midbrain nuclei\n"
        "HCPex_2mm.nii + HCPex_LookUpTable.txt  426 cortex/sub/midbrain/BF\n"
        "Tian_Subcortex_S1_3T.nii + labels      Official Tian-S1 volumetric\n"
        "FreeSurferColorLUT.txt                 ships with FreeSurfer install"
    )
    add_tree_box(s, MARGIN, row2_top + Emu(300000), quad_w, row2_h - Emu(300000), atlas_tree, size=8.3)

    # ---- Quadrant 4: Data / format prerequisites ----
    add_bullets_panel(
        s, quad2_left, row2_top, quad_w, row2_h,
        "Data & Format Prerequisites",
        [
            "Input already through HCP Pipelines/QuNex structural + functional "
            "preprocessing (PreFreeSurfer→FreeSurfer→PostFreeSurfer→fMRI Volume→"
            "Surface→ICA-FIX→MSMAll) — these scripts don't process raw DICOM.",
            "CIFTI grayordinate files: .dtseries.nii, .dscalar.nii, .dlabel.nii",
            "Volumetric NIfTI: .nii / .nii.gz",
            "Raw zip packages named <SubjectID>_<Visit>_MR_<Modality>Recommended*.zip",
        ],
        body_size=9.5,
    )

    # ================= HOW TO RUN =================
    s = new_content_slide(prs, blank_layout, "SUMMARY", "Correct Sequence — Running the Scripts")
    add_textbox(s, MARGIN, Emu(1230000), CONTENT_W, Emu(274320),
                "organize_hcp_data.py always first; the two analysis branches are independent of each other; documentation scripts run whenever, last.",
                11, SLATE_MUTED)

    col_gap = Emu(300000)
    col_w = (CONTENT_W - col_gap) // 2
    col2_left = MARGIN + col_w + col_gap

    add_textbox(s, MARGIN, Emu(1600000), col_w, Emu(280000),
                "STEP 0 — ONE-TIME SETUP", 13, NAVY, bold=True)
    add_textbox(s, col2_left, Emu(1600000), col_w, Emu(280000),
                "STEP 1 — organize_hcp_data.py (what it builds)", 13, NAVY, bold=True)

    step0_tree = (
        "atlases/  (must exist under Raw_Data/)\n"
        "├─ schaefer400_tianS1.dlabel.nii      Schaefer-400 + Tian-S1\n"
        "├─ CIT168_prob_func2mm.nii.gz+labels  5 midbrain nuclei\n"
        "├─ HCPex_2mm.nii + LookUpTable        426 cortex+sub+midbrain+BF\n"
        "└─ Tian_Subcortex_S1_3T.nii+labels    Official Tian-S1 volumetric\n"
        "\n"
        "FreeSurfer → FreeSurferColorLUT.txt\n"
        "FSL        → applywarp on PATH (cross-sectional v2 only)"
    )
    add_tree_box(s, MARGIN, Emu(1930000), col_w, Emu(1200000), step0_tree, size=9)

    step1_tree = (
        "Raw_Data/\n"
        "└─ sub-<subject>/\n"
        "    └─ ses-<visit>/\n"
        "        ├─ anat/     ← from *_StructuralRecommended*.zip\n"
        "        ├─ func/     ← from *_RestFmriRecommended*.zip\n"
        "        └─ concat/   ← from *_ConcatFmriRecommended*.zip"
    )
    add_tree_box(s, col2_left, Emu(1930000), col_w, Emu(1200000), step1_tree, size=9)

    add_textbox(s, MARGIN, Emu(3220000), CONTENT_W, Emu(280000),
                "STEP 2 — Two independent branches, either order, run only what you need", 14, NAVY, bold=True)
    rows_seq = [
        ["", "Branch A — Structural", "Branch B — Functional Connectivity"],
        ["Needs", "anat/  (from Step 1)", "func/  (from Step 1)  +  atlases"],
        ["Run", "cross_sectional_analysis_v1.py", "run_fc_pipeline_v2.py  (recommended, HCPex-only)  or  v1  (4-atlas legacy comparison)"],
        ["Recommended / alternative", "_v2.py — superset of v1, adds VTA/SN/NbM; "
         "or _batch_v2.py — every session unattended",
         "or the batch equivalent instead — run_fc_pipeline_batch_v2.py (v2 only) — every session unattended"],
    ]
    add_table(s, MARGIN, Emu(3560000), CONTENT_W,
              [Emu(2400000), Emu(4300000), Emu(4501400)],
              rows_seq, header_size=10.5, body_size=9.5, row_h=Emu(540000), highlight_rows={2})
    add_textbox(s, MARGIN, Emu(5760000), CONTENT_W, Emu(260000),
                "Branches don't depend on each other — they read different subfolders (anat/ vs func/) and write "
                "different subfolders of Analysed_data/<subject>/<session>/.",
                10, CAPTION_MUTED)

    add_textbox(s, MARGIN, Emu(6060000), CONTENT_W, Emu(300000),
                "STEP 3 (optional, anytime) — documentation/reporting scripts (build_scripts_overview_presentation.py, "
                "build_combined_presentation.py) just read already-computed results; not part of the data-processing path.",
                10, CAPTION_MUTED)

    # ================= CLOSING =================
    s = new_content_slide(prs, blank_layout, "SUMMARY", "Where to Go Next")
    add_bullets_panel(
        s, MARGIN, Emu(1400000), CONTENT_W, Emu(3200000),
        "Companion documents in data_analysis_scripts/",
        [
            "README.md — the narrative version of this deck (same content, same order).",
            "README.txt — one-line-per-script index: exact CLI flags, every output file "
            "path, per-script requirements. Kept up to date alongside every script.",
            "organize_hcp_data_flowchart.dot/.png/.svg, cross_sectional_analysis_flowchart.*, and "
            "fc_correlation_estimation_flowchart.* — standalone editable Graphviz sources for three "
            "of the four flowcharts shown in this deck. fc_pipeline_v1_flowchart.png (the fourth) "
            "has no editable source of its own.",
        ],
        body_size=13.5,
    )
    add_textbox(s, MARGIN, Emu(5100000), CONTENT_W, Emu(500000),
                "10 scripts total in this folder. Interactive scripts (run_fc_pipeline_v1/v2.py, "
                "cross_sectional_analysis_v1/v2.py) share the same subject/session tile-picker UI "
                "and each keep their own run-count log, so they never collide with each other's output.",
                11.5, SLATE_MUTED)

    prs.save(str(DST))
    print(f"Saved: {DST}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
